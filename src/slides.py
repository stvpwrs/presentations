"""Slide builder functions – one per layout type.

Each builder receives a ``Presentation``, slide data, a :class:`Style`, and
optional helpers for images and animations.
"""

from __future__ import annotations

import os
from typing import TYPE_CHECKING

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from pptx import Presentation

    from .style import Style

# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------


def _is_url(text: str) -> bool:
    """Return True if *text* is a bare URL."""
    return text.strip().startswith(("http://", "https://")) and " " not in text.strip()


def _apply_position(shape, pos: dict | None) -> None:
    """Move/resize *shape* according to a parsed position dict (inches)."""
    if not pos:
        return
    if "left" in pos:
        shape.left = Inches(pos["left"])
    if "top" in pos:
        shape.top = Inches(pos["top"])
    if "width" in pos:
        shape.width = Inches(pos["width"])
    if "height" in pos:
        shape.height = Inches(pos["height"])


def _add_image(slide, img: dict, pos: dict | None = None) -> None:
    """Add an image to *slide* from a parsed ``**Image**`` dict.

    If *pos* is given (from ``**ImagePos**``), it overrides the image dict coords.
    """
    path = img["path"]
    if not os.path.isfile(path):
        print(f"Warning: image not found: {path}, skipping.")
        return
    left = Inches((pos or {}).get("left", img.get("left", 6.5)))
    top = Inches((pos or {}).get("top", img.get("top", 1.5)))
    w_val = (pos or {}).get("width", img.get("width"))
    h_val = (pos or {}).get("height", img.get("height"))
    width = Inches(w_val) if w_val is not None else None
    height = Inches(h_val) if h_val is not None else None
    slide.shapes.add_picture(path, left, top, width, height)


# ---------------------------------------------------------------------------
# Layout builders
# ---------------------------------------------------------------------------


def _set_text_with_breaks(shape, text: str, font_size) -> None:
    """Set *text* on *shape*, splitting on ``<br>`` into separate paragraphs."""
    parts = [p.strip() for p in text.split("<br>")]
    tf = shape.text_frame
    tf.clear()
    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = part
        p.font.size = font_size


def add_title_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Layout 0 – Title Slide: large centred title + subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    positions = slide_data.get("positions", {})
    slide.shapes.title.text = slide_data["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.size = style.title_font
    subtitle = slide_data.get("subtitle", "")
    if "<br>" in subtitle:
        _set_text_with_breaks(slide.placeholders[1], subtitle, style.subtitle_font)
    else:
        slide.placeholders[1].text = subtitle
        slide.placeholders[1].text_frame.paragraphs[0].font.size = style.subtitle_font
    _apply_position(slide.shapes.title, positions.get("title"))
    _apply_position(slide.placeholders[1], positions.get("subtitle"))
    # Auto-align title top with image top when no explicit positions
    img = slide_data.get("image")
    if img and "title" not in positions and "image" not in positions:
        img_top = Inches(img.get("top", 1.3))
        slide.shapes.title.top = img_top
    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if slide_data.get("image"):
        _add_image(slide, slide_data["image"], positions.get("image"))
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


# Default positions (inches) for slide types without explicit Pos directives
_CONTENT_DEFAULTS = {
    "title": {"left": 0.5, "top": 0.2, "width": 8.5, "height": 0.8},
    "content": {"left": 0.5, "top": 1.2, "width": 6.5, "height": 5.5},
}
_TWO_COL_DEFAULTS = {
    "title": {"left": 0.5, "top": 0.2, "width": 8.5, "height": 0.8},
    "left": {"left": 0.5, "top": 1.2, "width": 4.25, "height": 4.0},
    "right": {"left": 5.0, "top": 1.2, "width": 4.25, "height": 4.0},
}
_SECTION_DEFAULTS = {
    "title": {"left": 0.5, "top": 2.0, "width": 9.0, "height": 1.5},
    "subtitle": {"left": 0.5, "top": 3.6, "width": 9.0, "height": 1.0},
}


def add_content_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Layout 1 – Title and Content: title bar + bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    positions = slide_data.get("positions", {})
    slide.shapes.title.text = slide_data["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.size = style.heading_font
    _apply_position(slide.shapes.title, positions.get("title") or _CONTENT_DEFAULTS["title"])
    body_ph = slide.shapes.placeholders[1]
    content_pos = positions.get("content") or _CONTENT_DEFAULTS["content"].copy()
    # Constrain content width when image is present to prevent overlap
    img = slide_data.get("image")
    if img and "content" not in positions:
        img_left = img.get("left", 6.5)
        content_pos["width"] = img_left - 0.2 - content_pos["left"]
    _apply_position(body_ph, content_pos)
    body = body_ph.text_frame
    body.clear()
    for i, b in enumerate(slide_data.get("bullets", [])):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = style.body_font
    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if slide_data.get("image"):
        _add_image(slide, slide_data["image"], positions.get("image"))
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


def add_section_header_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Layout 2 – Section Header: transition slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    positions = slide_data.get("positions", {})
    slide.shapes.title.text = slide_data["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.size = style.title_font
    _apply_position(slide.shapes.title, positions.get("title") or _SECTION_DEFAULTS["title"])
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        slide.placeholders[1].text = subtitle
        slide.placeholders[1].text_frame.paragraphs[0].font.size = style.subtitle_font
        _apply_position(slide.placeholders[1], positions.get("subtitle") or _SECTION_DEFAULTS["subtitle"])
    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if slide_data.get("image"):
        _add_image(slide, slide_data["image"], positions.get("image"))
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


def add_two_column_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Layout 3 – Two Content: side-by-side content placeholders."""
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    positions = slide_data.get("positions", {})
    slide.shapes.title.text = slide_data["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.size = style.heading_font
    _apply_position(slide.shapes.title, positions.get("title") or _TWO_COL_DEFAULTS["title"])

    left_placeholder = slide.placeholders[1]
    _apply_position(left_placeholder, positions.get("left") or _TWO_COL_DEFAULTS["left"])
    left_ph = left_placeholder.text_frame
    left_ph.clear()
    for i, b in enumerate(slide_data.get("left_bullets", [])):
        p = left_ph.paragraphs[0] if i == 0 else left_ph.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = style.col_body_font

    right_placeholder = slide.placeholders[2]
    _apply_position(right_placeholder, positions.get("right") or _TWO_COL_DEFAULTS["right"])
    right_ph = right_placeholder.text_frame
    right_ph.clear()
    for i, b in enumerate(slide_data.get("right_bullets", [])):
        p = right_ph.paragraphs[0] if i == 0 else right_ph.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = style.col_body_font

    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if slide_data.get("image"):
        _add_image(slide, slide_data["image"], positions.get("image"))
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' to an RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _interpolate_colors(colors: list[str], count: int) -> list[RGBColor]:
    """Distribute *colors* evenly across *count* slots via linear interpolation."""
    if count <= 0:
        return []
    if count == 1:
        return [_hex_to_rgb(colors[0])]
    result: list[RGBColor] = []
    for i in range(count):
        t = i / (count - 1)  # 0.0 → 1.0
        seg = t * (len(colors) - 1)
        lo = int(seg)
        hi = min(lo + 1, len(colors) - 1)
        frac = seg - lo
        c0 = _hex_to_rgb(colors[lo])
        c1 = _hex_to_rgb(colors[hi])
        r = int(c0[0] + (c1[0] - c0[0]) * frac)
        g = int(c0[1] + (c1[1] - c0[1]) * frac)
        b = int(c0[2] + (c1[2] - c0[2]) * frac)
        result.append(RGBColor(r, g, b))
    return result


def add_resource_box_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Custom layout – centred heading, coloured subtitle, labelled resource boxes."""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    import copy

    ss = slide_data.get("slide_style", {})

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide_w = prs.slide_width

    # --- Slide background ---
    bg_hex = ss.get("SlideBackground", style.slide_background)
    if bg_hex:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(bg_hex)

    # --- Title (h1) centred near top ---
    title_size = Pt(int(ss.get("TitleSize", "36")))
    title_color = _hex_to_rgb(ss["TitleColor"]) if ss.get("TitleColor") else RGBColor(0, 0, 0)

    title_left, title_top = Inches(0.5), Inches(0.4)
    title_width, title_height = slide_w - Inches(1.0), Inches(0.9)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = title_size
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.CENTER

    # --- Subtitle (h2) with per-character gradient colouring ---
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        sub_size = Pt(int(ss.get("SubtitleSize", "24")))
        sub_left, sub_top = Inches(0.5), Inches(1.3)
        sub_width, sub_height = slide_w - Inches(1.0), Inches(0.6)
        sub_box = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER

        gradient_raw = ss.get("SubtitleColors", "") or style.subtitle_colors
        if gradient_raw:
            color_stops = [c.strip() for c in gradient_raw.split(",")]
            char_colors = _interpolate_colors(color_stops, len(subtitle))
            sp.clear()
            for idx, ch in enumerate(subtitle):
                run = sp.add_run()
                run.text = ch
                run.font.size = sub_size
                run.font.color.rgb = char_colors[idx]
        else:
            sp.text = subtitle
            sp.font.size = sub_size
            sp.font.color.rgb = RGBColor(0x00, 0x78, 0xD4)

    # --- Resource boxes ---
    # Resolve from slide-level overrides (ss) falling back to global style
    badge_grad_start = ss.get("BadgeGradientStart", style.badge_gradient_start)
    badge_grad_end = ss.get("BadgeGradientEnd", style.badge_gradient_end)
    badge_text_color = _hex_to_rgb(ss.get("BadgeTextColor", style.badge_text_color))
    box_border_color = _hex_to_rgb(ss.get("BoxBorderColor", style.box_border_color))
    box_bg_hex = ss.get("BoxBackground", style.box_background)
    url_color = _hex_to_rgb(ss.get("UrlColor", style.url_color))
    name_color = _hex_to_rgb(ss.get("NameColor", style.name_color))
    divider_color = _hex_to_rgb(ss.get("DividerColor", style.divider_color))
    name_font_size = Pt(int(ss.get("NameFontSize", style.name_font_size)))
    url_font_size = Pt(int(ss.get("UrlFontSize", style.url_font_size)))
    badge_w = float(ss.get("BadgeWidth", style.badge_width))
    badge_h = float(ss.get("BadgeHeight", style.badge_height))
    badge_font_sz = int(ss.get("BadgeFontSize", style.badge_font_size))
    badge_corner = int(ss.get("BadgeCornerRadius", style.badge_corner_radius))
    box_corner = int(ss.get("BoxCornerRadius", style.box_corner_radius))

    box_top = Inches(2.3)
    outer_border_color_hex = ss.get("OuterBorderColor")
    for box_data in slide_data.get("boxes", []):
        label = box_data["label"]
        rows = box_data["rows"]

        # Container rounded rectangle (border + background)
        container_left = Inches(1.6)
        container_width = slide_w - Inches(2.0)
        num_rows = max(len(rows), 1)
        container_height = Inches(0.5 * num_rows + 0.35)

        # Optional outer border rectangle (drawn first so it sits behind)
        if outer_border_color_hex:
            outer_margin = Inches(0.12)
            badge_left_pos = Inches(0.3)
            outer_left = badge_left_pos - outer_margin
            outer_top = box_top - outer_margin
            outer_width = (container_left + container_width) - badge_left_pos + 2 * outer_margin
            outer_height = container_height + 2 * outer_margin
            outer = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                outer_left, outer_top, outer_width, outer_height,
            )
            outer.fill.background()  # transparent fill
            outer.line.color.rgb = _hex_to_rgb(outer_border_color_hex)
            outer.line.width = Pt(1.5)
            # Match corner rounding to container
            outer_sp = outer._element
            osp_pr = outer_sp.find(qn("p:spPr"))
            opr_elem = osp_pr.find(qn("a:prstGeom")) if osp_pr is not None else None
            if opr_elem is not None:
                oavLst = opr_elem.find(qn("a:avLst"))
                if oavLst is None:
                    oavLst = copy.deepcopy(opr_elem.makeelement(qn("a:avLst"), {}))
                    opr_elem.append(oavLst)
                oavLst.clear()
                ogd = oavLst.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {box_corner}"})
                oavLst.append(ogd)

        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            container_left, box_top, container_width, container_height,
        )
        container.fill.solid()
        container.fill.fore_color.rgb = _hex_to_rgb(box_bg_hex)
        container.line.color.rgb = box_border_color
        container.line.width = Pt(1.5)
        # Reduce corner rounding
        container_sp = container._element
        sp_pr_c = container_sp.find(qn("p:spPr"))
        pr_elem = sp_pr_c.find(qn("a:prstGeom")) if sp_pr_c is not None else None
        if pr_elem is not None:
            avLst = pr_elem.find(qn("a:avLst"))
            if avLst is None:
                avLst = copy.deepcopy(pr_elem.makeelement(qn("a:avLst"), {}))
                pr_elem.append(avLst)
            avLst.clear()
            gd = avLst.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {box_corner}"})
            avLst.append(gd)

        # Horizontal divider line inside container
        div_y = box_top + Inches(0.38)
        div_line = slide.shapes.add_connector(
            1, container_left + Inches(0.15), div_y,
            container_left + container_width - Inches(0.15), div_y,
        )
        div_line.line.color.rgb = divider_color
        div_line.line.width = Pt(0.75)

        # Resource text rows inside the container
        for i, row in enumerate(rows):
            row_y = box_top + Inches(0.08 + 0.5 * i)
            # Name text
            name_box = slide.shapes.add_textbox(
                container_left + Inches(0.2), row_y,
                Inches(2.5), Inches(0.35),
            )
            ntf = name_box.text_frame
            ntf.word_wrap = True
            np_ = ntf.paragraphs[0]
            np_.text = row["name"]
            np_.font.size = name_font_size
            np_.font.color.rgb = name_color

            # URL text (right-aligned)
            url_box = slide.shapes.add_textbox(
                container_left + container_width - Inches(4.7), row_y,
                Inches(4.5), Inches(0.35),
            )
            utf = url_box.text_frame
            utf.word_wrap = True
            up = utf.paragraphs[0]
            up.text = row["url"]
            up.font.size = url_font_size
            up.font.color.rgb = url_color
            up.alignment = PP_ALIGN.RIGHT

        # Badge / label shape (rounded rectangle with gradient)
        badge_left = Inches(0.3)
        badge_width, badge_height = Inches(badge_w), Inches(badge_h)
        badge_y = box_top + Inches((container_height / Inches(1) - badge_h) / 2)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            badge_left, badge_y, badge_width, badge_height,
        )
        # Gradient fill from start to end via XML
        # First clear any style-based fill by setting solid, then replacing with gradient
        badge.fill.solid()
        sp_pr = badge._element.find(qn("p:spPr"))
        # Remove the solidFill we just added
        for old_fill in list(sp_pr.findall(qn("a:solidFill"))):
            sp_pr.remove(old_fill)
        grad_fill = sp_pr.makeelement(qn("a:gradFill"), {})
        gs_lst = grad_fill.makeelement(qn("a:gsLst"), {})
        # Stop 1 – start color
        gs1 = gs_lst.makeelement(qn("a:gs"), {"pos": "0"})
        srgb1 = gs1.makeelement(qn("a:srgbClr"), {"val": badge_grad_start.lstrip("#")})
        gs1.append(srgb1)
        gs_lst.append(gs1)
        # Stop 2 – end color
        gs2 = gs_lst.makeelement(qn("a:gs"), {"pos": "100000"})
        srgb2 = gs2.makeelement(qn("a:srgbClr"), {"val": badge_grad_end.lstrip("#")})
        gs2.append(srgb2)
        gs_lst.append(gs2)
        grad_fill.append(gs_lst)
        # Linear gradient direction (top-to-bottom)
        lin = grad_fill.makeelement(qn("a:lin"), {"ang": "5400000", "scaled": "1"})
        grad_fill.append(lin)
        sp_pr.append(grad_fill)

        badge.line.fill.background()  # no border
        # Badge corner rounding
        badge_sp = badge._element
        badge_sp_pr = badge_sp.find(qn('p:spPr'))
        bpr = badge_sp_pr.find(qn('a:prstGeom')) if badge_sp_pr is not None else None
        if bpr is not None:
            bavLst = bpr.find(qn("a:avLst"))
            if bavLst is None:
                bavLst = copy.deepcopy(bpr.makeelement(qn("a:avLst"), {}))
                bpr.append(bavLst)
            bavLst.clear()
            bgd = bavLst.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {badge_corner}"})
            bavLst.append(bgd)

        btf = badge.text_frame
        btf.word_wrap = True
        from pptx.enum.text import MSO_ANCHOR
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.text = label
        bp.font.size = Pt(badge_font_sz)
        bp.font.color.rgb = badge_text_color
        bp.font.bold = True
        bp.alignment = PP_ALIGN.CENTER

        box_top += container_height + Inches(0.4)

    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


# ---------------------------------------------------------------------------
# Registry mapping slide type names → builder functions
# ---------------------------------------------------------------------------

from .rich_layouts import RICH_BUILDERS

SLIDE_BUILDERS: dict[str, callable] = {
    "title": add_title_slide,
    "content": add_content_slide,
    "section-header": add_section_header_slide,
    "two-column": add_two_column_slide,
    "resource-box": add_resource_box_slide,
    **RICH_BUILDERS,
}
