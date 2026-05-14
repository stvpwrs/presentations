"""Rich layout builders ported from the athenahealth dedicated script.

These builders consume a ``**Data**:`` YAML block on each slide (parsed into
``slide_data['data']``) plus optional ``Subtitle``/``Eyebrow``/``Footer``/
``Wordmark`` fields. They are intentionally opinionated about colors so the
generated decks have a consistent executive-report look out of the box.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
DARK = RGBColor(0x2D, 0x28, 0x95)
MID = RGBColor(0x6E, 0x69, 0xD9)
LIGHT = RGBColor(0xB8, 0xB5, 0xF5)
PALE = RGBColor(0xEE, 0xEA, 0xFB)
PALE_2 = RGBColor(0xE0, 0xDC, 0xF6)
GREEN = RGBColor(0x9A, 0xCA, 0x3C)
DOT_RED = RGBColor(0xE1, 0x1D, 0x48)
DOT_GREEN = RGBColor(0x16, 0xA3, 0x4A)
DOT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
SOFT_BG = RGBColor(0xF7, 0xF6, 0xFD)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Color & shape helpers
# ---------------------------------------------------------------------------
def _to_rgb(value, default: RGBColor) -> RGBColor:
    if isinstance(value, RGBColor):
        return value
    if isinstance(value, str) and value:
        return RGBColor.from_string(value.lstrip("#"))
    return default


_DOT_COLORS = {
    "red": DOT_RED, "green": DOT_GREEN, "amber": DOT_AMBER, "orange": DOT_AMBER,
    "yellow": DOT_AMBER, "gray": GRAY,
}

_PILL_COLORS = {
    "critical": DOT_RED, "urgent": DOT_AMBER, "high": MID,
    "medium": MID, "low": LIGHT_GRAY,
    "prod": DARK, "pre-prod": MID, "non-prod": LIGHT_GRAY,
}


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_oval(slide, x, y, w, h, fill, transparency=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    if transparency:
        solid = shp.fill._xPr.find(qn("a:solidFill"))
        if solid is not None:
            srgb = solid.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", str(int((100 - transparency) * 1000)))
    return shp


def add_rounded(slide, x, y, w, h, fill, line=None, radius=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI",
             italic=False, tracking=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if tracking:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(tracking))
    return tb


def _set_runs(tb, segments):
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for text, style in segments:
        run = p0.add_run()
        run.text = text
        run.font.name = style.get("font", "Segoe UI")
        run.font.size = Pt(style.get("size", 11))
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
        run.font.color.rgb = style.get("color", BLACK)


# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------
def _add_chrome(slide, wordmark="athenahealth"):
    if wordmark:
        add_text(slide, Inches(11.0), Inches(0.28), Inches(2.2), Inches(0.32),
                 wordmark, size=14, bold=True, color=DARK, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.5), Inches(7.18), Inches(12.33), Emu(20000),
             fill=PALE)


def _add_title(slide, title, subtitle):
    add_text(slide, Inches(0.5), Inches(0.45), Inches(11.5), Inches(0.6),
             title, size=30, bold=True, color=DARK)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=GRAY)


def _set_notes(slide, slide_data):
    notes = slide_data.get("notes", "") or ""
    slide.notes_slide.notes_text_frame.text = notes


# ---------------------------------------------------------------------------
# 1) hero-title
# ---------------------------------------------------------------------------
def add_hero_title_slide(prs, slide_data, style, *, apply_animations=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = slide_data.get("data") or {}
    bg = _to_rgb(data.get("background"), DARK)
    add_rect(slide, Emu(0), Emu(0), prs.slide_width, prs.slide_height, fill=bg)

    # Decorative circles on right
    add_oval(slide, Inches(8.2), Inches(-1.5), Inches(7.5), Inches(7.5),
             fill=MID, transparency=15)
    add_oval(slide, Inches(9.5), Inches(1.8), Inches(5.5), Inches(5.5),
             fill=LIGHT, transparency=10)
    add_oval(slide, Inches(10.5), Inches(3.5), Inches(3.5), Inches(3.5),
             fill=PALE, transparency=5)

    wordmark = slide_data.get("wordmark") or data.get("wordmark") or ""
    if wordmark:
        add_text(slide, Inches(0.6), Inches(0.5), Inches(4), Inches(0.4),
                 wordmark, size=20, bold=True, color=WHITE)

    eyebrow = slide_data.get("eyebrow") or data.get("eyebrow") or ""
    if eyebrow:
        add_text(slide, Inches(0.6), Inches(2.6), Inches(8), Inches(0.4),
                 eyebrow, size=12, bold=True, color=LIGHT, tracking=400)

    title_text = slide_data["title"].replace("<br>", "\n")
    add_text(slide, Inches(0.6), Inches(3.0), Inches(8.5), Inches(1.9),
             title_text, size=48, bold=True, color=WHITE)

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_text(slide, Inches(0.6), Inches(5.05), Inches(8), Inches(1.0),
                 subtitle.replace("<br>", "\n"), size=16, color=LIGHT)

    add_rect(slide, Inches(0.6), Inches(6.7), Inches(0.6), Emu(40000),
             fill=GREEN)

    footer = slide_data.get("footer") or data.get("footer") or ""
    if footer:
        add_text(slide, Inches(0.6), Inches(6.78), Inches(9), Inches(0.35),
                 footer, size=10, color=LIGHT)

    _set_notes(slide, slide_data)
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


# ---------------------------------------------------------------------------
# 2) stat-cards
# ---------------------------------------------------------------------------
def add_stat_cards_slide(prs, slide_data, style, *, apply_animations=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_chrome(slide, slide_data.get("wordmark") or "athenahealth")
    _add_title(slide, slide_data["title"], slide_data.get("subtitle", ""))

    data = slide_data.get("data") or {}
    cards = data.get("cards") or []
    n = max(len(cards), 1)
    card_w = Inches((12.33 - (n - 1) * 0.18) / n)
    card_h = Inches(2.55)
    gap = Inches(0.18)
    x0 = Inches(0.5)
    y0 = Inches(1.85)

    for i, card in enumerate(cards):
        x = x0 + i * (card_w + gap)
        add_rounded(slide, x, y0, card_w, card_h, fill=SOFT_BG, line=PALE, radius=0.04)
        add_rect(slide, x + Inches(0.25), y0 + Inches(0.3),
                 Inches(0.45), Emu(40000), fill=GREEN)
        add_text(slide, x + Inches(0.2), y0 + Inches(0.4),
                 card_w - Inches(0.4), Inches(1.2),
                 str(card.get("number", "")), size=44, bold=True, color=DARK)
        add_text(slide, x + Inches(0.25), y0 + Inches(1.65),
                 card_w - Inches(0.4), Inches(0.4),
                 card.get("label", ""), size=14, bold=True, color=DARK)
        if card.get("sub"):
            add_text(slide, x + Inches(0.25), y0 + Inches(2.05),
                     card_w - Inches(0.4), Inches(0.4),
                     card["sub"], size=11, color=GRAY)

    callouts = data.get("callouts") or []
    callout_title = data.get("callout_title", "")
    if callouts:
        why_y = Inches(4.7)
        why_h = Inches(2.0)
        add_rounded(slide, Inches(0.5), why_y, Inches(12.33), why_h,
                    fill=WHITE, line=PALE, radius=0.03)
        add_rect(slide, Inches(0.5), why_y, Inches(0.12), why_h, fill=GREEN)
        if callout_title:
            add_text(slide, Inches(0.85), Inches(4.85), Inches(8), Inches(0.4),
                     callout_title, size=16, bold=True, color=DARK)
        cn = len(callouts)
        col_w = Inches((12.33 - 0.35 - 0.1 * (cn - 1)) / cn)
        col_gap = Inches(0.1)
        for i, c in enumerate(callouts):
            cx = Inches(0.85) + i * (col_w + col_gap)
            add_text(slide, cx, Inches(5.3), col_w, Inches(0.35),
                     c.get("title", ""), size=12, bold=True, color=DARK)
            add_text(slide, cx, Inches(5.7), col_w, Inches(0.95),
                     c.get("body", ""), size=11, color=GRAY)

    _set_notes(slide, slide_data)


# ---------------------------------------------------------------------------
# 3) status-table — table with dot-status rows + optional side panel
# ---------------------------------------------------------------------------
def add_status_table_slide(prs, slide_data, style, *, apply_animations=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_chrome(slide, slide_data.get("wordmark") or "athenahealth")
    _add_title(slide, slide_data["title"], slide_data.get("subtitle", ""))

    data = slide_data.get("data") or {}
    rows = data.get("rows") or []
    columns = data.get("columns") or ["Model", "Deploys", "Status"]
    tbl_x = Inches(0.5)
    tbl_y = Inches(1.7)
    tbl_w = Inches(6.8)
    row_h = Inches(0.42)

    add_rect(slide, tbl_x, tbl_y, tbl_w, Inches(0.45), fill=DARK)
    col_xs = [Inches(0.15), Inches(3.35), Inches(4.55)]
    col_ws = [Inches(2.9), Inches(0.9), Inches(2.1)]
    col_aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT]
    for c, lbl in enumerate(columns[:3]):
        add_text(slide, tbl_x + col_xs[c], tbl_y, col_ws[c], Inches(0.45),
                 lbl, size=12, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=col_aligns[c])

    for i, row in enumerate(rows):
        ry = tbl_y + Inches(0.45) + i * row_h
        fill = PALE if i % 2 == 0 else WHITE
        add_rect(slide, tbl_x, ry, tbl_w, row_h, fill=fill)
        add_text(slide, tbl_x + Inches(0.15), ry, Inches(2.9), row_h,
                 str(row.get("model", row.get("name", ""))),
                 size=12, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tbl_x + Inches(3.35), ry, Inches(0.9), row_h,
                 str(row.get("count", row.get("deploys", ""))),
                 size=14, bold=True, color=DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        dot_color = _DOT_COLORS.get((row.get("dot") or "").lower(), GRAY)
        dot_size = Inches(0.16)
        add_oval(slide, tbl_x + Inches(4.55), ry + (row_h - dot_size) / 2,
                 dot_size, dot_size, fill=dot_color)
        add_text(slide, tbl_x + Inches(4.78), ry, Inches(2), row_h,
                 str(row.get("status", "")), size=11, color=GRAY,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Right panel — bar list
    rx = Inches(7.7)
    side = data.get("side") or {}
    side_title = side.get("title", "")
    bars = side.get("bars") or []
    if side_title:
        add_text(slide, rx, Inches(1.7), Inches(5.1), Inches(0.35),
                 side_title, size=14, bold=True, color=DARK)
    bar_colors = [DARK, MID, LIGHT, PALE_2]
    for i, b in enumerate(bars):
        y = Inches(2.2) + i * Inches(0.58)
        add_text(slide, rx, y, Inches(3), Inches(0.3),
                 b.get("name", ""), size=12, bold=True, color=DARK)
        pct = float(b.get("pct", 0))
        add_text(slide, rx + Inches(4.0), y, Inches(1.1), Inches(0.3),
                 f"{b.get('count','')}  ({pct:g}%)", size=11, bold=True,
                 color=DARK, align=PP_ALIGN.RIGHT)
        add_rect(slide, rx, y + Inches(0.32), Inches(5.1), Inches(0.16),
                 fill=PALE_2)
        fw = Inches(5.1 * pct / 100.0)
        add_rect(slide, rx, y + Inches(0.32), fw, Inches(0.16),
                 fill=bar_colors[i % len(bar_colors)])

    callout = data.get("callout") or {}
    if callout:
        rrx = Inches(7.7)
        rry = Inches(4.05)
        rrw = Inches(5.13)
        rrh = Inches(2.6)
        add_rounded(slide, rrx, rry, rrw, rrh, fill=DARK, radius=0.05)
        if callout.get("eyebrow"):
            add_text(slide, rrx + Inches(0.3), rry + Inches(0.15),
                     Inches(4.5), Inches(0.3), callout["eyebrow"],
                     size=11, bold=True, color=LIGHT, tracking=300)
        if callout.get("number"):
            add_text(slide, rrx + Inches(0.3), rry + Inches(0.45),
                     Inches(4.5), Inches(0.85), callout["number"],
                     size=44, bold=True, color=WHITE)
        if callout.get("body"):
            add_text(slide, rrx + Inches(0.3), rry + Inches(1.45),
                     Inches(4.55), Inches(0.7), callout["body"],
                     size=12, color=LIGHT)
        add_rect(slide, rrx + Inches(0.3), rry + Inches(2.0),
                 Inches(0.4), Emu(30000), fill=GREEN)
        if callout.get("path"):
            tb = add_text(slide, rrx + Inches(0.3), rry + Inches(2.1),
                          Inches(4.6), Inches(0.4), "Path",
                          size=11, italic=True, color=LIGHT)
            _set_runs(tb, [
                ("Path: ", {"italic": True, "color": LIGHT, "size": 11}),
                (callout["path"],
                 {"italic": True, "bold": True, "color": WHITE, "size": 11}),
            ])

    if data.get("caption"):
        add_text(slide, Inches(0.5), Inches(6.63), Inches(8), Inches(0.4),
                 data["caption"], size=10, italic=True, color=GRAY)

    _set_notes(slide, slide_data)


# ---------------------------------------------------------------------------
# 4) stack-table — env pill + stacked bar + total per row
# ---------------------------------------------------------------------------
def add_stack_table_slide(prs, slide_data, style, *, apply_animations=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_chrome(slide, slide_data.get("wordmark") or "athenahealth")
    _add_title(slide, slide_data["title"], slide_data.get("subtitle", ""))

    data = slide_data.get("data") or {}
    rows = data.get("rows") or []
    segments_def = data.get("segments") or [
        {"key": "ptu", "label": "PTU", "color": "#2D2895"},
        {"key": "standard", "label": "Standard", "color": "#6E69D9"},
        {"key": "other", "label": "Other", "color": "#B8B5F5"},
    ]
    seg_colors = [_to_rgb(s.get("color"), DARK) for s in segments_def]
    seg_keys = [s["key"] for s in segments_def]

    cols_y = Inches(1.7)
    add_text(slide, Inches(0.5), cols_y, Inches(2.7), Inches(0.3),
             "Subscription", size=11, bold=True, color=GRAY, tracking=200)
    add_text(slide, Inches(3.2), cols_y, Inches(1.6), Inches(0.3),
             "Environment", size=11, bold=True, color=GRAY, tracking=200)
    add_text(slide, Inches(4.6), cols_y, Inches(6), Inches(0.3),
             "   /   ".join(s["label"] for s in segments_def),
             size=11, bold=True, color=GRAY, tracking=200)
    add_text(slide, Inches(10.75), cols_y, Inches(0.7), Inches(0.3),
             "Total", size=11, bold=True, color=GRAY,
             align=PP_ALIGN.RIGHT, tracking=200)

    chart_left = Inches(4.6)
    chart_max_w = 6.0
    max_total = max((sum(r.get(k, 0) for k in seg_keys) for r in rows), default=1)
    if max_total == 0:
        max_total = 1

    for i, r in enumerate(rows):
        ry = Inches(2.1) + i * Inches(0.55)
        add_text(slide, Inches(0.5), ry, Inches(2.7), Inches(0.4),
                 r.get("name", ""), size=12, bold=True, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        env = r.get("env", "")
        pill_w = Inches(0.95)
        pill_h = Inches(0.26)
        py = ry + Inches(0.07)
        pill_color = _PILL_COLORS.get(env.lower(), GRAY)
        add_rounded(slide, Inches(3.2), py, pill_w, pill_h,
                    fill=pill_color, radius=0.5)
        add_text(slide, Inches(3.2), py, pill_w, pill_h,
                 env, size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        seg_y = ry + Inches(0.08)
        seg_h = Inches(0.26)
        x = chart_left
        scale = chart_max_w / max_total
        for k, color in zip(seg_keys, seg_colors):
            count = r.get(k, 0)
            if count > 0:
                w = Inches(count * scale)
                add_rect(slide, x, seg_y, w, seg_h, fill=color)
                x += w
        total = r.get("total", sum(r.get(k, 0) for k in seg_keys))
        add_text(slide, Inches(10.75), ry, Inches(0.7), Inches(0.4),
                 str(total), size=14, bold=True, color=DARK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Legend
    leg_y = Inches(6.0)
    legend_xs = [4.6, 6.6, 8.6, 10.6]
    for i, seg in enumerate(segments_def):
        x = legend_xs[min(i, len(legend_xs) - 1)]
        add_rect(slide, Inches(x), leg_y + Inches(0.08), Inches(0.2),
                 Inches(0.18), fill=seg_colors[i])
        label = seg.get("legend_label", seg["label"])
        add_text(slide, Inches(x + 0.28), leg_y, Inches(2.4),
                 Inches(0.35), label, size=11, color=GRAY)

    callout = data.get("callout") or {}
    if callout:
        fy = Inches(6.45)
        fh = Inches(0.6)
        add_rounded(slide, Inches(0.5), fy, Inches(12.33), fh,
                    fill=SOFT_BG, line=PALE, radius=0.2)
        add_rect(slide, Inches(0.5), fy, Inches(0.12), fh, fill=GREEN)
        add_text(slide, Inches(0.8), fy + Inches(0.05), Inches(2.6),
                 Inches(0.3), callout.get("title", ""), size=12,
                 bold=True, color=DARK)
        add_text(slide, Inches(3.4), fy, Inches(9.3), fh,
                 callout.get("body", ""), size=11, color=GRAY,
                 anchor=MSO_ANCHOR.MIDDLE)

    _set_notes(slide, slide_data)


# ---------------------------------------------------------------------------
# 5) priority-table — rows with priority pill badges
# ---------------------------------------------------------------------------
def add_priority_table_slide(prs, slide_data, style, *, apply_animations=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_chrome(slide, slide_data.get("wordmark") or "athenahealth")
    _add_title(slide, slide_data["title"], slide_data.get("subtitle", ""))

    data = slide_data.get("data") or {}
    rows = data.get("rows") or []
    columns = data.get("columns") or ["Priority", "Action", "Replacement", "Deploys", "Deadline"]

    hy = Inches(1.7)
    add_rect(slide, Inches(0.5), hy, Inches(12.33), Inches(0.5), fill=DARK)
    layout = [
        (0.65, 1.2,  PP_ALIGN.LEFT),
        (2.15, 4.2,  PP_ALIGN.LEFT),
        (6.65, 3.2,  PP_ALIGN.LEFT),
        (10.15, 1.1, PP_ALIGN.CENTER),
        (11.55, 1.13,PP_ALIGN.LEFT),
    ]
    for (x, w, align), lbl in zip(layout, columns):
        add_text(slide, Inches(x), hy, Inches(w), Inches(0.5),
                 lbl, size=11, bold=True, color=WHITE,
                 align=align, anchor=MSO_ANCHOR.MIDDLE, tracking=200)

    rh = Inches(0.46)
    for i, r in enumerate(rows):
        ry = Inches(2.2) + i * rh
        fill = PALE if i % 2 == 0 else WHITE
        add_rect(slide, Inches(0.5), ry, Inches(12.33), rh, fill=fill)
        py = ry + Inches(0.09)
        priority = r.get("priority", "")
        pill_color = _to_rgb(
            r.get("color"),
            _PILL_COLORS.get(priority.lower(), MID),
        )
        add_rounded(slide, Inches(0.65), py, Inches(1.2), Inches(0.28),
                    fill=pill_color, radius=0.5)
        add_text(slide, Inches(0.65), py, Inches(1.2), Inches(0.28),
                 priority, size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=300)
        add_text(slide, Inches(2.15), ry, Inches(4.2), rh,
                 r.get("action", ""), size=12, bold=True, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(6.65), ry, Inches(3.2), rh,
                 r.get("replacement", ""), size=12, color=GRAY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(10.15), ry, Inches(1.1), rh,
                 str(r.get("deploys", "")), size=14, bold=True, color=DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(11.55), ry, Inches(1.4), rh,
                 r.get("deadline", ""), size=11, color=GRAY,
                 anchor=MSO_ANCHOR.MIDDLE)

    _set_notes(slide, slide_data)


# ---------------------------------------------------------------------------
# 6) timeline-cards — wave cards along a horizontal rail
# ---------------------------------------------------------------------------
def add_timeline_cards_slide(prs, slide_data, style, *, apply_animations=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_chrome(slide, slide_data.get("wordmark") or "athenahealth")
    _add_title(slide, slide_data["title"], slide_data.get("subtitle", ""))

    data = slide_data.get("data") or {}
    cards = data.get("cards") or []
    n = max(len(cards), 1)
    card_w = Inches((12.33 - (n - 1) * 0.13) / n)
    card_h = Inches(4.4)
    card_y = Inches(3.0)
    gap = Inches(0.13)
    x0 = Inches(0.5)
    line_y = Inches(2.53)
    add_rect(slide, Inches(0.7), line_y, Inches(11.93), Emu(50000), fill=PALE)

    for i, c in enumerate(cards):
        cx = x0 + i * (card_w + gap)
        dot_size = Inches(0.3)
        dot_cx = cx + card_w / 2 - dot_size / 2
        add_oval(slide, dot_cx, line_y + Inches(0.03) - dot_size / 2 + Inches(0.025),
                 dot_size, dot_size, fill=DARK)
        add_rounded(slide, cx, card_y, card_w, card_h, fill=SOFT_BG,
                    line=PALE, radius=0.04)
        add_rect(slide, cx, card_y, card_w, Inches(0.12), fill=DARK)
        add_text(slide, cx + Inches(0.25), card_y + Inches(0.25),
                 card_w - Inches(0.5), Inches(0.3),
                 c.get("label", ""), size=11, bold=True, color=DARK,
                 tracking=300)
        add_text(slide, cx + Inches(0.25), card_y + Inches(0.55),
                 card_w - Inches(0.5), Inches(0.5),
                 c.get("phase", ""), size=22, bold=True, color=DARK)
        add_text(slide, cx + Inches(0.25), card_y + Inches(1.05),
                 card_w - Inches(0.5), Inches(0.3),
                 c.get("when", ""), size=11, color=GRAY)
        add_text(slide, cx + Inches(0.25), card_y + Inches(1.4),
                 card_w - Inches(0.5), Inches(0.3),
                 c.get("count", ""), size=12, bold=True, color=DARK)
        add_rect(slide, cx + Inches(0.25), card_y + Inches(1.78),
                 card_w - Inches(0.5), Emu(20000), fill=PALE)
        for j, item in enumerate(c.get("items") or []):
            add_text(slide, cx + Inches(0.25),
                     card_y + Inches(1.9) + j * Inches(0.45),
                     card_w - Inches(0.5), Inches(0.4),
                     item, size=11, color=DARK)

    _set_notes(slide, slide_data)


# ---------------------------------------------------------------------------
# Registry
# ---------------------------------------------------------------------------
RICH_BUILDERS = {
    "hero-title":     add_hero_title_slide,
    "stat-cards":     add_stat_cards_slide,
    "status-table":   add_status_table_slide,
    "stack-table":    add_stack_table_slide,
    "priority-table": add_priority_table_slide,
    "timeline-cards": add_timeline_cards_slide,
}
