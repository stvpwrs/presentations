# /// script
# requires-python = ">=3.10"
# dependencies = ["ms-presentations==0.1.6"]
# ///
"""Self-contained MRE — Bug 5: 4:3-vs-16:9 geometry mismatch for core layouts.

Run with:  uv run mse.py   (no venv, no Azure, no network)

Bug:      A rich layout flips the deck to 16:9, but core builders use 10"-canvas
          default positions, so core slides don't span the widescreen canvas.
Expected: On a 13.333" canvas the core title/body scale to fill the slide.
Actual:   The content slide's title stays 8.5" wide on a 13.333" canvas (an
          off-center title + ~6" right gap). The printed widths below are the
          deterministic observable.
"""
import os
import tempfile

from pptx import Presentation
from presentations.cli import main

SPEC = """\
---
title: Repro - Aspect ratio mismatch
output: Bug05_Aspect_Ratio_Mismatch.pptx
---

## [hero-title] Rich Layout Forces Widescreen

**Eyebrow**: KEYNOTE
**Subtitle**: This rich slide flips the deck to 16:9 (13.333" wide).

**Notes**: hero-title is a rich type, so the whole deck renders widescreen.

---

## [content] Core Slide On The Wrong Canvas

- Core content slide with default geometry.
- Its title default is 8.5" wide, sized for a 10" canvas, not 13.333".

**Notes**: On the forced widescreen canvas the title no longer spans the slide.
"""

with tempfile.TemporaryDirectory() as d:
    spec_path = os.path.join(d, "minimal.spec.md")
    with open(spec_path, "w", encoding="utf-8") as f:
        f.write(SPEC)
    main([spec_path, "-o", d])

    prs = Presentation(os.path.join(d, "Bug05_Aspect_Ratio_Mismatch.pptx"))
    print("slide width (in):", round(prs.slide_width / 914400, 3))  # ~13.333
    content_slide = prs.slides[1]
    if content_slide.shapes.title is not None:
        print("core title width (in):", round(content_slide.shapes.title.width / 914400, 3))
    print("Title width ~8.5 on a ~13.333 canvas -> does not span the slide.")
